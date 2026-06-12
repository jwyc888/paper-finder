"""
paperfinder.py: Tier A core.

Pipeline:  capture -> queue -> metadata pass (instant) -> embed pass (background)
           -> hybrid index (FTS5 keyword + dense vectors, fused by RRF).

Staged ingestion is the point:
  * metadata pass  = cheap. Title + first-page text indexed to FTS5 immediately,
                     so a just-dropped paper is keyword-findable within seconds.
  * embed pass     = heavier. Full-text parse + embedding, run in the background;
                     adds semantic recall and full-text keyword reach.

doc_id is the CANONICAL identity, shared with the relationship layer. Embeddings
live in their own column and are a regenerable cache, and re-embedding never touches
identity or any verified relationship.

Dense search here is brute-force cosine (fine at personal scale, zero extra deps);
swap in sqlite-vec when the corpus grows. The embedder defaults to a dependency-free
hashing embedder; point it at bge-small / PubMedBERT for real semantic recall.
"""

from __future__ import annotations

import io
import json
import hashlib
import math
import os
import re
import sqlite3
import time
from typing import Optional

import pypdf

from paperfinder.core.capture import CaptureSource, DocumentRef
from paperfinder.core.vectorstore import BruteForceStore, VectorStore

DOI_RE = re.compile(r"10\.\d{4,9}/[-._;()/:A-Z0-9]+", re.I)


# --------------------------------------------------------------------------- #
# Parsing
# --------------------------------------------------------------------------- #
class Parser:
    @staticmethod
    def quick(ref: DocumentRef, data: bytes) -> dict:
        """Cheap metadata-pass extraction: title, first-page text, doi."""
        if ref.kind == "pdf":
            try:
                reader = pypdf.PdfReader(io.BytesIO(data))
                page0 = (reader.pages[0].extract_text() or "") if reader.pages else ""
                meta_title = (reader.metadata.title if reader.metadata else None) or ""
            except Exception:
                page0, meta_title = "", ""
            title = meta_title.strip() or _first_line(page0) or ref.name
            return {"title": title, "first_text": page0[:2000],
                    "doi": _doi(page0), "kind": "pdf"}
        if ref.kind in ("text", "docx", "pptx"):
            txt = Parser.full(ref, data)  # these formats are small; read once
            return {"title": _first_line(txt) or ref.name,
                    "first_text": txt[:2000], "doi": _doi(txt), "kind": ref.kind}
        # url / other: nothing to parse cheaply
        return {"title": ref.name, "first_text": "", "doi": None, "kind": ref.kind}

    @staticmethod
    def full(ref: DocumentRef, data: bytes) -> str:
        """Embed-pass extraction: full text."""
        if ref.kind == "pdf":
            try:
                reader = pypdf.PdfReader(io.BytesIO(data))
                return "\n".join((p.extract_text() or "") for p in reader.pages)
            except Exception:
                return ""
        if ref.kind == "text":
            return _decode(data)
        if ref.kind == "docx":
            return _docx_text(data)
        if ref.kind == "pptx":
            return _pptx_text(data)
        return ""


def _docx_text(data: bytes) -> str:
    """Text from a .docx (paragraphs + table cells). Graceful if the optional
    `python-docx` dependency is missing (install the `office` extra)."""
    try:
        import docx
    except ImportError:
        return ""
    try:
        d = docx.Document(io.BytesIO(data))
        parts = [p.text for p in d.paragraphs]
        for table in d.tables:
            for row in table.rows:
                parts.extend(cell.text for cell in row.cells)
        return "\n".join(t for t in parts if t)
    except Exception:
        return ""


def _pptx_text(data: bytes) -> str:
    """Text from a .pptx (all text-bearing shapes across slides). Graceful if the
    optional `python-pptx` dependency is missing."""
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    try:
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)
    except Exception:
        return ""


def _decode(data: bytes) -> str:
    try:
        return data.decode("utf-8", "replace")
    except Exception:
        return ""


def _first_line(text: str) -> str:
    for line in text.splitlines():
        s = line.strip()
        if len(s) > 8:
            return s[:160]
    return ""


def _doi(text: str) -> Optional[str]:
    m = DOI_RE.search(text or "")
    return m.group(0) if m else None


def _chunk_text(text: str, size: int = 350, overlap: int = 50) -> list[str]:
    """Split text into overlapping ~`size`-word passages. Word-based (not token-
    based) so it's tokenizer-agnostic; ~350 words stays comfortably under a 512-token
    model limit. Overlap keeps ideas that straddle a boundary findable in both."""
    words = (text or "").split()
    if not words:
        return []
    if len(words) <= size:
        return [" ".join(words)]
    step = max(1, size - overlap)
    chunks = []
    for start in range(0, len(words), step):
        window = words[start:start + size]
        if not window:
            break
        chunks.append(" ".join(window))
        if start + size >= len(words):
            break
    return chunks


def _snippet(text: str, terms: list[str], width: int = 300) -> str:
    """A passage preview centred on the first matching query term, so the evidence
    shows *why* the chunk matched rather than just its opening words. Falls back to
    the head when there's no lexical hit (e.g. a purely semantic match)."""
    low = (text or "").lower()
    pos = -1
    for t in terms:
        i = low.find(t)
        if i != -1 and (pos == -1 or i < pos):
            pos = i
    if pos == -1:
        return text[:width] + ("…" if len(text) > width else "")
    start = max(0, pos - width // 3)
    end = min(len(text), start + width)
    return ("…" if start > 0 else "") + text[start:end] + ("…" if end < len(text) else "")


# --------------------------------------------------------------------------- #
# Embedding (pluggable; default is dependency-free)
# --------------------------------------------------------------------------- #
class HashingEmbedder:
    """Deterministic hashed bag-of-words, L2-normalised. A stand-in so the
    pipeline runs with zero model downloads. Lexical, not semantic, so swap for a
    real model on the Mac for genuine recall."""
    model_name = "hashing-v0"

    def __init__(self, dim: int = 256):
        self.dim = dim

    def embed(self, text: str) -> list[float]:
        vec = [0.0] * self.dim
        for tok in re.findall(r"[a-z0-9]+", (text or "").lower()):
            # stable across processes (built-in hash() is salted by PYTHONHASHSEED,
            # which would change vectors between the embed run and a later query)
            h = int.from_bytes(hashlib.blake2b(tok.encode(), digest_size=8).digest(), "big")
            vec[h % self.dim] += 1.0
        n = math.sqrt(sum(v * v for v in vec))
        return [v / n for v in vec] if n else vec


# To use real semantic embeddings on the Mac, select this via PAPERFINDER_EMBEDDER=st.
class STEmbedder:
    """Real semantic embeddings via sentence-transformers (e.g. bge-small).
    Lazy import so the package is only required if you actually select it.
    UNTESTED in this sandbox (no model downloaded), so verify on the Mac."""

    def __init__(self, model: str = "BAAI/bge-small-en-v1.5"):
        from sentence_transformers import SentenceTransformer
        self.model_name = model
        self._m = SentenceTransformer(model)

    def embed(self, text: str) -> list[float]:
        return self._m.encode(text or "", normalize_embeddings=True).tolist()


def cosine(u: list[float], v: list[float]) -> float:
    # kept for backwards-compat imports; dense scoring now lives in the VectorStore.
    from paperfinder.core.vectorstore import cosine as _c
    return _c(u, v)


# --------------------------------------------------------------------------- #
# Core
# --------------------------------------------------------------------------- #
class PaperFinder:
    def __init__(self, db_path: str = "paperfinder.db", embedder=None,
                 vector_store=None, vector_store_name: str = "bruteforce"):
        self.conn = sqlite3.connect(db_path, check_same_thread=False)
        self.conn.row_factory = sqlite3.Row
        self.embedder = embedder or HashingEmbedder()
        # opt-in: drop references/bibliography/supplementary before chunking
        self.strip_sections = bool(os.environ.get("PAPERFINDER_STRIP_SECTIONS"))
        # opt-in: section-aware chunking (splits on detected sections, tags each
        # chunk, drops back-matter). Default off keeps legacy flat chunking.
        self.section_chunks = bool(os.environ.get("PAPERFINDER_SECTION_CHUNKS"))
        self._segment_classify = None   # injectable for tests (else the local LLM)
        self._init_schema()
        # the pluggable dense backend, pass an instance or name one to build
        if vector_store is not None:
            self.store: VectorStore = vector_store
        else:
            from paperfinder.core.vectorstore import make_store
            dim = len(self.embedder.embed("dimension probe"))
            self.store = make_store(vector_store_name, self.conn, dim)

    def _init_schema(self) -> None:
        self.conn.executescript(
            """
            CREATE TABLE IF NOT EXISTS documents (
                doc_id TEXT PRIMARY KEY,
                title TEXT, source_url TEXT, kind TEXT, doi TEXT,
                descriptors TEXT,                 -- JSON; filled by humans/LLM later
                first_text TEXT, full_text TEXT,
                embedding_model TEXT,
                modified REAL, indexed_at REAL, embedded_at REAL,
                folder TEXT,                      -- full path from the source root ("" = root)
                archived INTEGER DEFAULT 0
            );
            CREATE VIRTUAL TABLE IF NOT EXISTS docs_fts
                USING fts5(doc_id UNINDEXED, content);
            CREATE TABLE IF NOT EXISTS chunks (
                chunk_id TEXT PRIMARY KEY,
                doc_id TEXT, ordinal INTEGER, text TEXT,
                section_text TEXT, section_type TEXT
            );
            CREATE INDEX IF NOT EXISTS idx_chunks_doc ON chunks(doc_id);
            CREATE TABLE IF NOT EXISTS jobs (
                job_id INTEGER PRIMARY KEY AUTOINCREMENT,
                doc_id TEXT, ref TEXT, stage TEXT, status TEXT, created_at REAL
            );
            CREATE TABLE IF NOT EXISTS meta (key TEXT PRIMARY KEY, value TEXT);
            CREATE INDEX IF NOT EXISTS idx_jobs_status ON jobs(stage, status);
            """
        )
        for ddl in (  # migrate pre-existing DBs (each is a no-op once the column exists)
            "ALTER TABLE documents ADD COLUMN archived INTEGER DEFAULT 0",
            "ALTER TABLE documents ADD COLUMN folder TEXT",
            "ALTER TABLE chunks ADD COLUMN section_text TEXT",
            "ALTER TABLE chunks ADD COLUMN section_type TEXT",
        ):
            try:
                self.conn.execute(ddl)
            except sqlite3.OperationalError:
                pass
        self.conn.commit()

    # ---- checkpoint ------------------------------------------------------
    def _get_meta(self, key: str) -> Optional[str]:
        r = self.conn.execute("SELECT value FROM meta WHERE key=?", (key,)).fetchone()
        return r["value"] if r else None

    def _set_meta(self, key: str, value: Optional[str]) -> None:
        self.conn.execute(
            "INSERT INTO meta(key,value) VALUES(?,?) "
            "ON CONFLICT(key) DO UPDATE SET value=excluded.value", (key, value))
        self.conn.commit()

    # ---- passage building (section-aware optional) -----------------------
    def _build_passages(self, title: str, body: str, fallback: str) -> list[dict]:
        """Return passages to embed, each {text, section_text, section_type}.

        Section-aware path (PAPERFINDER_SECTION_CHUNKS): split on detected
        sections, drop back-matter, window WITHIN each section (no cross-section
        overlap), and tag every chunk with the section's verbatim heading and
        normalized type. Flat path (default): the original word-window over
        title+body with no section labels, so legacy behaviour is unchanged."""
        combined = f"{title}\n{body}"
        if self.section_chunks:
            from paperfinder.core.sectionstrip import segment, DROP_TYPES
            spans = (segment(combined, classify=self._segment_classify)
                     if self._segment_classify is not None else segment(combined))
            lines = combined.splitlines()
            out = []
            for sp in spans:
                if sp["section_type"] in DROP_TYPES:        # references/supplementary
                    continue
                span_text = "\n".join(lines[sp["start"]:sp["end"]])
                for w in _chunk_text(span_text):
                    out.append({"text": w, "section_text": sp["section_text"],
                                "section_type": sp["section_type"]})
            if out:
                return out
            return [{"text": fallback, "section_text": "", "section_type": "other"}]
        passages = _chunk_text(combined) or [fallback]
        return [{"text": w, "section_text": "", "section_type": None} for w in passages]

    # ---- capture ---------------------------------------------------------
    def run_capture(self, source: CaptureSource, source_key: str = "default") -> int:
        ckpt = self._get_meta(f"ckpt:{source_key}")
        refs, new_ckpt = source.poll(ckpt)
        for ref in refs:
            self.conn.execute(
                "INSERT INTO jobs(doc_id,ref,stage,status,created_at) "
                "VALUES(?,?,?,?,?)",
                (ref.doc_id, _ref_json(ref), "metadata", "pending", time.time()))
        self._set_meta(f"ckpt:{source_key}", new_ckpt)
        self.conn.commit()
        # stash refs so passes can fetch bytes without re-polling
        self._refs = getattr(self, "_refs", {})
        for ref in refs:
            self._refs[ref.doc_id] = ref
        return len(refs)

    def run_backfill(self, source, source_key: str = "default",
                     reconcile: bool = False) -> dict:
        """One-time, in-place backfill. crawl() if available (Drive, follows
        aliases), else a full local re-scan. Enqueues only new/returning docs.
        With reconcile=True, papers no longer reachable get archived (see reconcile)."""
        self._refs = getattr(self, "_refs", {})
        if hasattr(source, "crawl"):
            refs = source.crawl()
            prefix, new_ckpt = "gdrive:", source.start_checkpoint()
        else:
            self._set_meta(f"ckpt:{source_key}", None)        # force full re-scan
            refs, new_ckpt = source.poll(None)
            prefix = "local:"
        reachable = set()
        enqueued = 0
        for ref in refs:
            self._refs[ref.doc_id] = ref
            reachable.add(ref.doc_id)
            existing = self.get_document(ref.doc_id)
            if existing and not existing["archived"]:
                new_tag = getattr(ref, "tag", "") or ""
                if (existing["folder"] or "") != new_tag:     # folder renamed, or doc moved
                    self.conn.execute("UPDATE documents SET folder=? WHERE doc_id=?",
                                      (new_tag, ref.doc_id))
                continue                                       # already in scope; poll handles content edits
            self.conn.execute(
                "INSERT INTO jobs(doc_id,ref,stage,status,created_at) VALUES(?,?,?,?,?)",
                (ref.doc_id, _ref_json(ref), "metadata", "pending", time.time()))
            enqueued += 1
        self._set_meta(f"ckpt:{source_key}", new_ckpt)
        self.conn.commit()
        archived = self.reconcile(reachable, prefix) if reconcile else 0
        return {"reachable": len(reachable), "enqueued": enqueued, "archived": archived}

    def reconcile(self, reachable_ids, source_prefix: str) -> int:
        """Soft-archive indexed docs from this source that are no longer reachable
        (a paper deleted, or a folder/alias removed from scope). Only the LOCAL
        index is touched, never the source. Archived docs drop out of search; the
        document row and any authenticated relationships are preserved, so it's
        reversible: re-indexing the same doc un-archives it."""
        reachable = set(reachable_ids)
        archived = 0
        for d in self.all_documents():
            if not d["doc_id"].startswith(source_prefix):
                continue
            if d["doc_id"] in reachable or d["archived"]:
                continue
            self.conn.execute("UPDATE documents SET archived=1 WHERE doc_id=?", (d["doc_id"],))
            self.conn.execute("DELETE FROM docs_fts WHERE doc_id=?", (d["doc_id"],))
            self._delete_chunks(d["doc_id"])
            archived += 1
        self.conn.commit()
        return archived

    def _rehydrate(self, job: sqlite3.Row) -> DocumentRef:
        """Reconstruct a fetchable ref for a queued job. Uses the in-process ref
        if present; otherwise rebuilds it from the durable job record (local
        files), so the passes are re-runnable after a crash or across calls."""
        meta = json.loads(job["ref"])
        if meta["doc_id"] in getattr(self, "_refs", {}):
            return self._refs[meta["doc_id"]]
        url = meta.get("source_url", "")
        if url.startswith("file://"):
            path = url[len("file://"):]

            def _fetch(p=path) -> bytes:
                with open(p, "rb") as f:
                    return f.read()

            return DocumentRef(meta["doc_id"], meta["name"], meta["kind"],
                               meta.get("modified", 0.0), url, _fetch)
        raise RuntimeError(
            f"no fetcher for {meta['doc_id']}: non-local sources must be polled in-process")

    # ---- staged passes ---------------------------------------------------
    def run_metadata_pass(self) -> int:
        rows = self.conn.execute(
            "SELECT * FROM jobs WHERE stage='metadata' AND status='pending'").fetchall()
        done = 0
        for job in rows:
            ref = self._rehydrate(job)
            data = ref.fetch()
            m = Parser.quick(ref, data)
            now = time.time()
            self.conn.execute(
                """INSERT INTO documents
                   (doc_id,title,source_url,kind,doi,descriptors,first_text,folder,modified,indexed_at,archived)
                   VALUES(?,?,?,?,?,?,?,?,?,?,0)
                   ON CONFLICT(doc_id) DO UPDATE SET
                     title=excluded.title, source_url=excluded.source_url,
                     kind=excluded.kind, doi=excluded.doi,
                     first_text=excluded.first_text, folder=excluded.folder,
                     indexed_at=excluded.indexed_at,
                     archived=0""",
                (ref.doc_id, m["title"], ref.source_url, m["kind"], m["doi"],
                 json.dumps([]), m["first_text"], getattr(ref, "tag", ""), ref.modified, now))
            self._fts_set(ref.doc_id, f"{m['title']}\n{m['first_text']}")
            self.conn.execute("UPDATE jobs SET status='done' WHERE job_id=?", (job["job_id"],))
            self.conn.execute(
                "INSERT INTO jobs(doc_id,ref,stage,status,created_at) VALUES(?,?,?,?,?)",
                (ref.doc_id, job["ref"], "embed", "pending", now))
            done += 1
        self.conn.commit()
        return done

    def run_embed_pass(self) -> int:
        rows = self.conn.execute(
            "SELECT * FROM jobs WHERE stage='embed' AND status='pending'").fetchall()
        done = 0
        for job in rows:
            ref = self._rehydrate(job)
            data = ref.fetch()
            full = Parser.full(ref, data)
            if self.strip_sections and not self.section_chunks and full:
                from paperfinder.core.sectionstrip import strip_back_matter
                full = strip_back_matter(full)
            doc = self.get_document(ref.doc_id)
            base = full or doc["first_text"] or ""
            passages = self._build_passages(doc["title"] or "", base, doc["title"] or ref.name)
            now = time.time()
            self._delete_chunks(ref.doc_id)            # re-embed safety: replace cleanly
            for i, p in enumerate(passages):
                cid = f"{ref.doc_id}#{i}"
                self.store.upsert(cid, self.embedder.embed(p["text"]))
                self.conn.execute(
                    "INSERT INTO chunks(chunk_id,doc_id,ordinal,text,section_text,section_type) "
                    "VALUES(?,?,?,?,?,?)",
                    (cid, ref.doc_id, i, p["text"], p["section_text"], p["section_type"]))
            self.conn.execute(
                """UPDATE documents SET full_text=?, embedding_model=?, embedded_at=?
                   WHERE doc_id=?""",
                (full, self.embedder.model_name, now, ref.doc_id))
            # widen keyword reach to the full text now that we have it
            self._fts_set(ref.doc_id, f"{doc['title']}\n{base}")
            self.conn.execute("UPDATE jobs SET status='done' WHERE job_id=?", (job["job_id"],))
            done += 1
        self.conn.commit()
        return done

    def reembed_all(self, embedder) -> None:
        """Re-embed every document with a new model, re-chunking from full text.
        Identity + FTS keys are unchanged, so nothing downstream (relationships) is
        disturbed. Note: a model with a different vector dimension needs a fresh
        store (brute-force tolerates it; sqlite-vec/Qdrant fix dim at creation)."""
        self.embedder = embedder
        for d in self.all_documents():
            base = d["full_text"] or d["first_text"] or ""
            if self.strip_sections and not self.section_chunks and base:
                from paperfinder.core.sectionstrip import strip_back_matter
                base = strip_back_matter(base)
            passages = self._build_passages(d["title"] or "", base, d["title"] or "")
            self._delete_chunks(d["doc_id"])
            for i, p in enumerate(passages):
                cid = f"{d['doc_id']}#{i}"
                self.store.upsert(cid, embedder.embed(p["text"]))
                self.conn.execute(
                    "INSERT INTO chunks(chunk_id,doc_id,ordinal,text,section_text,section_type) "
                    "VALUES(?,?,?,?,?,?)",
                    (cid, d["doc_id"], i, p["text"], p["section_text"], p["section_type"]))
            self.conn.execute(
                "UPDATE documents SET full_text=?, embedding_model=?, embedded_at=? WHERE doc_id=?",
                (base, embedder.model_name, time.time(), d["doc_id"]))
        self.conn.commit()

    # ---- chunk helpers ---------------------------------------------------
    def _delete_chunks(self, doc_id: str) -> None:
        for r in self.conn.execute("SELECT chunk_id FROM chunks WHERE doc_id=?", (doc_id,)):
            self.store.delete(r["chunk_id"])
        self.conn.execute("DELETE FROM chunks WHERE doc_id=?", (doc_id,))

    def doc_vector(self, doc_id: str) -> Optional[list[float]]:
        """Centroid of a document's chunk vectors, a single-vector view of a doc
        for callers (e.g. the relationship layer) that still want one per document."""
        vecs = []
        for r in self.conn.execute("SELECT chunk_id FROM chunks WHERE doc_id=?", (doc_id,)):
            v = self.store.get(r["chunk_id"])
            if v:
                vecs.append(v)
        if not vecs:
            return None
        dim = len(vecs[0])
        cent = [sum(v[i] for v in vecs) / len(vecs) for i in range(dim)]
        n = math.sqrt(sum(x * x for x in cent))
        return [x / n for x in cent] if n else cent

    def add_document_text(self, doc_id: str, title: str, text: str,
                          source_url: str = "", doi: Optional[str] = None,
                          kind: str = "text", folder: str = "") -> None:
        """Index a document directly from in-memory text (metadata + chunked embed
        in one call). Convenience for programmatic ingestion and tests."""
        if self.strip_sections and not self.section_chunks and text:
            from paperfinder.core.sectionstrip import strip_back_matter
            text = strip_back_matter(text)
        now = time.time()
        self.conn.execute(
            """INSERT INTO documents
               (doc_id,title,source_url,kind,doi,descriptors,first_text,full_text,folder,
                embedding_model,modified,indexed_at,embedded_at,archived)
               VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,0)
               ON CONFLICT(doc_id) DO UPDATE SET
                 title=excluded.title, source_url=excluded.source_url, kind=excluded.kind,
                 doi=excluded.doi, first_text=excluded.first_text, full_text=excluded.full_text,
                 folder=excluded.folder,
                 embedding_model=excluded.embedding_model, indexed_at=excluded.indexed_at,
                 embedded_at=excluded.embedded_at, archived=0""",
            (doc_id, title, source_url, kind, doi, json.dumps([]),
             text[:2000], text, folder, self.embedder.model_name, now, now, now))
        self._fts_set(doc_id, f"{title}\n{text}")
        passages = self._build_passages(title or "", text, title or doc_id)
        self._delete_chunks(doc_id)
        for i, p in enumerate(passages):
            cid = f"{doc_id}#{i}"
            self.store.upsert(cid, self.embedder.embed(p["text"]))
            self.conn.execute(
                "INSERT INTO chunks(chunk_id,doc_id,ordinal,text,section_text,section_type) "
                "VALUES(?,?,?,?,?,?)",
                (cid, doc_id, i, p["text"], p["section_text"], p["section_type"]))
        self.conn.commit()

    # ---- fts helper ------------------------------------------------------
    def _fts_set(self, doc_id: str, content: str) -> None:
        self.conn.execute("DELETE FROM docs_fts WHERE doc_id=?", (doc_id,))
        self.conn.execute("INSERT INTO docs_fts(doc_id,content) VALUES(?,?)", (doc_id, content))

    # ---- reads -----------------------------------------------------------
    def get_document(self, doc_id: str) -> Optional[dict]:
        r = self.conn.execute("SELECT * FROM documents WHERE doc_id=?", (doc_id,)).fetchone()
        return dict(r) if r else None

    def all_documents(self) -> list[dict]:
        return [dict(r) for r in self.conn.execute("SELECT * FROM documents")]

    # ---- connection discovery (chunk-neighbour candidate edges) ----------
    def propose_connections(self, doc_id: str, k: int = 5, per_chunk: int = 8,
                            min_score: float = 0.0) -> list[dict]:
        """Find documents related to `doc_id` by NEAREST PASSAGE, not whole-doc average.
        For each of this doc's chunks, pull nearby chunks, drop same-doc hits, and keep
        each other document's single best passage pair. Returns ranked candidates, each
        carrying the two passages that justify the connection (the evidence)."""
        src_chunks = self.conn.execute(
            "SELECT chunk_id, text FROM chunks WHERE doc_id=?", (doc_id,)).fetchall()
        best: dict[str, dict] = {}
        for sc in src_chunks:
            v = self.store.get(sc["chunk_id"])
            if not v:
                continue
            for cid, score in self.store.query(v, per_chunk + 1):
                if score < min_score:
                    continue
                r = self.conn.execute(
                    "SELECT doc_id, text FROM chunks WHERE chunk_id=?", (cid,)).fetchone()
                if not r or r["doc_id"] == doc_id:
                    continue
                od = r["doc_id"]
                if od not in best or score > best[od]["score"]:
                    best[od] = {"score": score,
                                "src_passage": sc["text"], "dst_passage": r["text"]}
        ranked = sorted(best.items(), key=lambda x: x[1]["score"], reverse=True)[:k]
        out = []
        for od, info in ranked:
            d = self.get_document(od)
            out.append({
                "doc_id": od,
                "title": d["title"] if d else od,
                "score": round(info["score"], 5),
                "src_passage": info["src_passage"],   # full matching passages = the evidence
                "dst_passage": info["dst_passage"],
            })
        return out

    def build_graph_candidates(self, rg, k: int = 5, min_score: float = 0.0) -> int:
        """Populate a RelationshipGraph with chunk-neighbour candidate edges across the
        whole (active) corpus. `rg` is duck-typed (add_document / record_candidate), so
        the finder stays decoupled from the graph module. Returns edges proposed."""
        proposed = 0
        for d in self.all_documents():
            if d["archived"]:
                continue
            rg.add_document(d["doc_id"], d["title"], [], self.doc_vector(d["doc_id"]) or [],
                            source_url=d["source_url"], embedding_model=self.embedder.model_name,
                            folder=d["folder"] or "")
        for d in self.all_documents():
            if d["archived"]:
                continue
            for c in self.propose_connections(d["doc_id"], k=k, min_score=min_score):
                rg.record_candidate(
                    d["doc_id"], c["doc_id"], score=c["score"],
                    evidence={"src_passage": c["src_passage"], "dst_passage": c["dst_passage"]})
                proposed += 1
        return proposed

    # ---- hybrid search ---------------------------------------------------
    def search(self, query: str, k: int = 5, rrf_k: int = 60,
               folder: Optional[str] = None, section: Optional[str] = None) -> list[dict]:
        # keyword ranker (BM25 via FTS5)
        terms = re.findall(r"[a-z0-9]+", query.lower())
        kw_rank: dict[str, int] = {}
        if terms:
            match = " OR ".join(f'"{t}"' for t in terms)
            try:
                rows = self.conn.execute(
                    "SELECT doc_id FROM docs_fts WHERE docs_fts MATCH ? "
                    "ORDER BY bm25(docs_fts) LIMIT 50", (match,)).fetchall()
                for i, r in enumerate(rows):
                    kw_rank[r["doc_id"]] = i
            except sqlite3.OperationalError:
                pass

        # dense ranker: nearest CHUNKS, pooled up to their documents (best passage wins)
        qv = self.embedder.embed(query)
        doc_best: dict[str, tuple] = {}   # doc_id -> (score, chunk_id, section_text, section_type)
        for cid, score in self.store.query(qv, 50):
            row = self.conn.execute(
                "SELECT doc_id, section_text, section_type FROM chunks WHERE chunk_id=?",
                (cid,)).fetchone()
            if not row:
                continue
            if section and (row["section_type"] or "") != section:   # section-scoped retrieval
                continue
            did = row["doc_id"]
            if did not in doc_best or score > doc_best[did][0]:
                doc_best[did] = (score, cid, row["section_text"], row["section_type"])
        dense_sorted = sorted(doc_best.items(), key=lambda x: x[1][0], reverse=True)
        dense_rank = {did: i for i, (did, _) in enumerate(dense_sorted)}

        # reciprocal-rank fusion
        fused: dict[str, float] = {}
        for ranks in (kw_rank, dense_rank):
            for doc_id, rank in ranks.items():
                fused[doc_id] = fused.get(doc_id, 0.0) + 1.0 / (rrf_k + rank)

        out = []
        for doc_id, score in sorted(fused.items(), key=lambda x: x[1], reverse=True):
            d = self.get_document(doc_id)
            if not d or d["archived"]:
                continue
            if folder:                                # exact folder or anything beneath it
                df = d.get("folder") or ""
                if df != folder and not df.startswith(folder + "/"):
                    continue
            if section and doc_id not in doc_best:    # section-scoped: only docs with a matching span
                continue
            passage = None
            section_text = ""
            section_type = None
            if doc_id in doc_best:
                section_text = doc_best[doc_id][2] or ""
                section_type = doc_best[doc_id][3]
                pr = self.conn.execute(
                    "SELECT text FROM chunks WHERE chunk_id=?", (doc_best[doc_id][1],)).fetchone()
                if pr and pr["text"]:
                    passage = _snippet(pr["text"], terms)
            out.append({
                "doc_id": doc_id,                     # canonical identity
                "title": d["title"],
                "source_url": d["source_url"],        # the link a human re-opens
                "folder": d.get("folder") or "",      # full path from the source root
                "doi": d["doi"],
                "embedded": d["embedding_model"] is not None,
                "passage": passage,                   # the matching passage (the "why")
                "section_text": section_text,         # verbatim heading of the matching chunk
                "section_type": section_type,         # normalized bucket (None on flat-chunked docs)
                "score": round(score, 5),
            })
            if len(out) >= k:
                break
        return out


def _ref_json(ref: DocumentRef) -> str:
    return json.dumps({"doc_id": ref.doc_id, "name": ref.name, "kind": ref.kind,
                       "source_url": ref.source_url, "modified": ref.modified})

"""Verify Word (.docx) and PowerPoint (.pptx) get text-extracted and indexed.

Needs the office extra:  pip install -e ".[office]"
Run:  python3 tests/test_filetypes.py
"""

import os
import shutil
import sys
import tempfile

from docx import Document
from pptx import Presentation

from paperfinder.core.capture import LocalFolderSource
from paperfinder.core.finder import HashingEmbedder, PaperFinder

TOPIC = "patient sentiment toward AI chatbots in clinical care"


def main() -> int:
    d = tempfile.mkdtemp()

    doc = Document()
    doc.add_paragraph("Review of " + TOPIC)
    doc.add_paragraph("This Word document discusses " + TOPIC + " in depth.")
    doc.save(os.path.join(d, "review.docx"))

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slides on " + TOPIC
    prs.save(os.path.join(d, "deck.pptx"))

    db = "test_filetypes.db"
    if os.path.exists(db):
        os.remove(db)
    pf = PaperFinder(db, embedder=HashingEmbedder())
    pf.run_backfill(LocalFolderSource(d), source_key="local", reconcile=True)
    pf.run_metadata_pass()
    pf.run_embed_pass()

    kinds = {os.path.basename(x["source_url"]): x["kind"] for x in pf.all_documents()}
    hits = {os.path.basename(pf.get_document(h["doc_id"])["source_url"])
            for h in pf.search(TOPIC, k=10)}

    shutil.rmtree(d, ignore_errors=True)
    checks = [
        ("docx recognized as docx", kinds.get("review.docx") == "docx"),
        ("pptx recognized as pptx", kinds.get("deck.pptx") == "pptx"),
        ("docx searchable by its content", "review.docx" in hits),
        ("pptx searchable by its content", "deck.pptx" in hits),
    ]
    ok = True
    for name, passed in checks:
        print(f"  [{'PASS' if passed else 'FAIL'}] {name}")
        ok = ok and passed
    print("\n" + ("ALL CHECKS PASSED" if ok else "SOME CHECKS FAILED"))
    return 0 if ok else 1


if __name__ == "__main__":
    sys.exit(main())
